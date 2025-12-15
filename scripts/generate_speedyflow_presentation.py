from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()

PRIMARY = RGBColor(99, 102, 241)   # #6366f1
ACCENT = RGBColor(6, 182, 212)     # #06b6d4
TEXT_DARK = RGBColor(31, 41, 55)   # #1f2937
MUTED = RGBColor(156, 163, 175)    # #9ca3af

slides_content = [
    {
        'title': 'SpeedyFlow vs Jira Web',
        'bullets': ['Acelerando la gestión de incidencias: velocidad, contexto y automatización'],
        'notes': 'Apertura: presentar objetivo de la comparación - productividad, costes y UX.'
    },
    {
        'title': 'Problema actual con Jira Web',
        'bullets': [
            'Interfaz genérica → fricción en flujos de soporte',
            'Contexto fragmentado; cambio de pestañas y pérdida de hilo',
            'Automatizaciones limitadas; tareas repetitivas',
            'Rendimiento en vistas con muchos tickets (latencia)'
        ],
        'notes': 'Establecer dolores: espera, contexto perdido, trabajo repetitivo.'
    },
    {
        'title': 'Propuesta de valor de SpeedyFlow',
        'bullets': [
            'Interfaz optimizada para soporte: tableros y atajos contextuales',
            'Cachés y hashing inteligente → renders ultrarrápidos',
            'AI/ML integrado: sugerencias, duplicados, prioridad predictiva',
            'Acciones en contexto sin salir del flujo'
        ],
        'notes': 'Enfatizar ahorro de tiempo y mejor experiencia de agentes.'
    },
    {
        'title': 'Beneficios cuantificables',
        'bullets': [
            'Reducción tiempo medio por ticket: -30% a -60% (estimado)',
            'Menos re-trabajo/duplicados: hasta -40%',
            'Interacciones rápidas: vistas cacheadas <100ms',
            'Mejora en CSAT y cumplimiento de SLA'
        ],
        'notes': 'Presentar como estimaciones conservadoras; proponer piloto para medición real.'
    },
    {
        'title': 'Killer features de SpeedyFlow',
        'bullets': [
            'Triple capa de caching (sidebar, tablero, issue)',
            'FlowingContext: sugerencias contextuales en tiempo real',
            'Hash-based rendering: re-render sólo si cambia el ticket',
            'Integraciones AI plug-and-play y privacidad local',
            'Kanban/Queue híbrido con reglas automáticas'
        ],
        'notes': 'Breve explicación de cada feature; por qué es diferencial frente a Jira Web.'
    },
    {
        'title': 'Caso de uso: resultado esperado',
        'bullets': [
            'Escenario: 50 agentes, 5k tickets/mes',
            'Piloto 4 semanas → Tiempo de resolución: -35%',
            'Tickets duplicados/reasignados: -30%',
            'Carga de tableros: -80%'
        ],
        'notes': 'Proponer piloto controlado y KPIs (TTR, duplicados, CSAT, latencia).' 
    },
    {
        'title': 'Próximos pasos & CTA',
        'bullets': [
            'Proponer piloto 4 semanas (2 equipos) y medir KPIs',
            'Entregables: integración mínima, formación 2h, reporte de impacto',
            'Agendar demo técnica y roadmap de integración'
        ],
        'notes': 'Cerrar con CTA claro: agendar demo y plan de despliegue.'
    }
]

def set_title_style(title_shape):
    tf = title_shape.text_frame
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = PRIMARY


def add_icon_box(slide, left, top, emoji, bg_color):
    w = Inches(0.6)
    h = Inches(0.6)
    box = slide.shapes.add_shape(1, left, top, w, h)  # MSO_SHAPE_RECTANGLE (value 1)
    # Fill with color
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    box.line.fill.background()
    # add emoji text
    tx = box.text_frame
    tx.text = emoji
    p = tx.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255,255,255)


for idx, s in enumerate(slides_content):
    slide_layout = prs.slide_layouts[1]  # title and content
    slide = prs.slides.add_slide(slide_layout)

    # Title and small icon accent
    title = slide.shapes.title
    title.text = s['title']
    set_title_style(title)
    add_icon_box(slide, Inches(0.3), Inches(0.4), '💡' if idx==0 else '⚡', PRIMARY)

    # content placeholder
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, bullet in enumerate(s['bullets']):
        p = tf.add_paragraph() if i>0 or len(tf.paragraphs)==0 else tf.paragraphs[0]
        # prefix icon for bullets
        icon = ''
        if 'Killer' in s['title']:
            icon = '🛠️ '
        elif 'Beneficios' in s['title'] or 'Beneficios cuantificables' in s['title']:
            icon = '📈 '
        elif 'Problema' in s['title']:
            icon = '⚠️ '
        elif 'Próximos' in s['title']:
            icon = '✅ '
        elif 'Caso de uso' in s['title']:
            icon = '🔍 '
        else:
            icon = '🔹 '
        p.text = f"{icon}{bullet}"
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_DARK

    # speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = s.get('notes', '')

# Title slide (first) customization: make bigger hero
if prs.slides:
    first = prs.slides[0]
    first.shapes.title.text = 'SpeedyFlow vs Jira Web'
    set_title_style(first.shapes.title)

# Save output
import os
os.makedirs('output', exist_ok=True)
from datetime import datetime
ts = datetime.now().strftime('%Y%m%d_%H%M%S')
outfile = os.path.join('output', f'SpeedyFlow_vs_Jira_{ts}.pptx')
prs.save(outfile)
print('Created', outfile)
