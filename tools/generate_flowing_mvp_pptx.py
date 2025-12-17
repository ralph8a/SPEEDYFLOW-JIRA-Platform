from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pathlib import Path
OUT = 'output/Flowing_MVP_Capabilities.pptx'
def set_title(slide, text):
    title = slide.shapes.title
    title.text = text
    for paragraph in title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Calibri'
            run.font.size = Pt(32)
            run.font.bold = True
            run.font.color.rgb = RGBColor(37,99,235)
def add_bullet_slide(prs, title_text, bullets):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    set_title(slide, title_text)
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        p = body.add_paragraph() if i>0 else body.paragraphs[0]
        p.text = b
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(18)
            run.font.name = 'Calibri'
            run.font.color.rgb = RGBColor(31,41,55)
    return slide
def style_slide_background(slide):
    # subtle gradient-like solid color to match brand
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(249,250,251)  # very light
def main():
    prs = Presentation()
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = 'Flowing MVP — Capacidades'
    slide.placeholders[1].text = 'Contextual AI assistant for tickets, suggestions and automation'
    style_slide_background(slide)
    add_bullet_slide(prs, 'Asistente contextual inteligente', [
        'Detecta contexto (kanban, lista, tarjeta, sidebar)',
        'Sugerencias ML: resumir, sugerir respuestas, autocompletar campos',
        'Búsqueda semántica y detección de duplicados',
    ])
    add_bullet_slide(prs, 'Automatización y análisis ML', [
        'AI Queue Analyzer: insights, clustering y anomalías',
        'Reglas/blueprints ML para asignaciones y prioridades',
        'Módulos desacoplados: frontend ml-client y backend blueprints'
    ])
    add_bullet_slide(prs, 'Integración y UX', [
        'Editor inline asistido por AI y plantillas guardables',
        'Integración vía blueprints backend (`/api/flowing/*`)',
        'Control de privacidad y modo debug/opt-in'
    ])
    # Apply background style to all slides
    for s in prs.slides:
        style_slide_background(s)
    out_path = Path(OUT)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f'Written PPTX to {out_path}')
if __name__ == '__main__':
    main()
